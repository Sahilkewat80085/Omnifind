from datetime import datetime, timezone
from pathlib import Path
from pptx import Presentation

from core.document_pipeline.extractors.base import BaseExtractor
from core.document_pipeline.models import (
    FormatMetadata,
    NormalizedBlock,
    NormalizedDocument,
    NormalizedTable,
)


class PptxExtractor(BaseExtractor):
    """Extractor for PowerPoint (.pptx) presentations extracting slides, speaker notes, and tables."""

    def extract(self, file_path: str) -> NormalizedDocument:
        path = Path(file_path)
        stat = path.stat()
        created_at = datetime.fromtimestamp(stat.st_ctime, tz=timezone.utc)
        modified_at = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)

        prs = Presentation(str(path))
        blocks: list[NormalizedBlock] = []
        tables: list[NormalizedTable] = []
        slide_count = len(prs.slides)
        has_speaker_notes = False

        table_counter = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_loc = f"Slide {slide_idx}"
            slide_title = ""

            # Check for slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
                blocks.append(
                    NormalizedBlock(
                        text=slide_title,
                        location=slide_loc,
                        heading_path=[slide_title],
                        block_type="heading",
                        level=1,
                    )
                )

            # Process shapes in slide
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                # 1. Text frames / bullet points
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        blocks.append(
                            NormalizedBlock(
                                text=text,
                                location=slide_loc,
                                heading_path=[slide_title] if slide_title else [],
                                block_type="paragraph",
                            )
                        )

                # 2. Slide Tables
                if shape.has_table:
                    table_counter += 1
                    t = shape.table
                    table_rows: list[list[str]] = []
                    for row in t.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_rows.append(row_data)

                    if table_rows:
                        headers = table_rows[0]
                        data_rows = table_rows[1:] if len(table_rows) > 1 else []
                        tables.append(
                            NormalizedTable(
                                table_id=f"table-{table_counter}",
                                location=slide_loc,
                                headers=headers,
                                rows=data_rows,
                            )
                        )

            # 3. Speaker Notes Extraction
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_speaker_notes = True
                    blocks.append(
                        NormalizedBlock(
                            text=f"Speaker Notes: {notes_text}",
                            location=f"{slide_loc} (Notes)",
                            heading_path=[slide_title] if slide_title else [],
                            block_type="speaker_note",
                        )
                    )

        return NormalizedDocument(
            file_path=str(path.resolve()),
            file_name=path.name,
            file_type="pptx",
            created_at=created_at,
            modified_at=modified_at,
            blocks=blocks,
            tables=tables,
            format_metadata=FormatMetadata(
                slide_count=slide_count,
                has_speaker_notes=has_speaker_notes,
            ),
            is_scanned=False,
            requires_ocr=False,
        )
