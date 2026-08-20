from pathlib import Path
import docx
import fitz
import openpyxl
from pptx import Presentation

from core.document_pipeline.extractors.docx_extractor import DocxExtractor
from core.document_pipeline.extractors.pdf_extractor import PdfExtractor
from core.document_pipeline.extractors.pptx_extractor import PptxExtractor
from core.document_pipeline.extractors.txt_md_extractor import TxtMdExtractor
from core.document_pipeline.extractors.xlsx_extractor import XlsxExtractor
from core.document_pipeline.pipeline import DocumentPipeline
from core.document_pipeline.router import FormatRouter
from core.document_pipeline.semantic.chunker import SemanticChunker
from core.document_pipeline.semantic.keywords import KeywordExtractor
from core.document_pipeline.semantic.ner import EntityRecognizer


def test_txt_md_extractor(tmp_path):
    md_file = tmp_path / "sample.md"
    md_file.write_text(
        "# Project Overview\n\n"
        "This is a high-level summary of the architecture.\n\n"
        "## Components\n\n"
        "- Frontend: React\n"
        "- Backend: FastAPI\n\n"
        "| Service | Port |\n"
        "|---|---|\n"
        "| API | 8000 |\n"
        "| UI | 5173 |\n",
        encoding="utf-8",
    )

    extractor = TxtMdExtractor()
    doc = extractor.extract(str(md_file))

    assert doc.file_type == "md"
    assert len(doc.blocks) >= 3
    assert len(doc.tables) == 1
    assert doc.tables[0].headers == ["Service", "Port"]
    assert doc.tables[0].rows == [["API", "8000"], ["UI", "5173"]]


def test_docx_extractor(tmp_path):
    docx_file = tmp_path / "test.docx"
    doc_obj = docx.Document()
    doc_obj.add_heading("Quarterly Earnings", level=1)
    doc_obj.add_paragraph("Total revenue grew by 25% year over year.")
    table = doc_obj.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Revenue"
    table.rows[1].cells[1].text = "$10M"
    doc_obj.save(str(docx_file))

    extractor = DocxExtractor()
    doc = extractor.extract(str(docx_file))

    assert doc.file_type == "docx"
    assert any(b.text == "Quarterly Earnings" and b.block_type == "heading" for b in doc.blocks)
    assert len(doc.tables) == 1
    assert doc.tables[0].headers == ["Metric", "Value"]


def test_pptx_extractor(tmp_path):
    pptx_file = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "OmniFind Pitch"
    slide.placeholders[1].text = "AI-Powered Context Search"
    slide.notes_slide.notes_text_frame.text = "Emphasize local-first zero-cost operation."
    prs.save(str(pptx_file))

    extractor = PptxExtractor()
    doc = extractor.extract(str(pptx_file))

    assert doc.file_type == "pptx"
    assert doc.format_metadata.slide_count == 1
    assert doc.format_metadata.has_speaker_notes is True
    assert any("Emphasize local-first" in b.text for b in doc.blocks)


def test_xlsx_extractor(tmp_path):
    xlsx_file = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "SalesData"
    ws.append(["Region", "Sales", "Rep"])
    ws.append(["North", 50000, "Alice"])
    ws.append(["South", 75000, "Bob"])
    wb.save(str(xlsx_file))

    extractor = XlsxExtractor()
    doc = extractor.extract(str(xlsx_file))

    assert doc.file_type == "xlsx"
    assert doc.format_metadata.sheet_names == ["SalesData"]
    assert len(doc.tables) == 1
    assert doc.tables[0].headers == ["Region", "Sales", "Rep"]


def test_scanned_pdf_detection(tmp_path):
    pdf_file = tmp_path / "empty_scan.pdf"
    doc_obj = fitz.open()
    # Add an empty page with no text
    doc_obj.new_page()
    doc_obj.save(str(pdf_file))
    doc_obj.close()

    extractor = PdfExtractor()
    doc = extractor.extract(str(pdf_file))

    assert doc.file_type == "pdf"
    assert doc.is_scanned is True
    assert doc.requires_ocr is True


def test_format_router(tmp_path):
    router = FormatRouter()
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Meeting notes from yesterday.", encoding="utf-8")

    extractor = router.get_extractor(str(txt_file))
    assert isinstance(extractor, TxtMdExtractor)


def test_semantic_chunker_preserves_locations(tmp_path):
    md_file = tmp_path / "doc.md"
    md_file.write_text("# Title\n\n" + "Word " * 200 + "\n\n# Section 2\n\n" + "Data " * 200, encoding="utf-8")
    doc = TxtMdExtractor().extract(str(md_file))

    chunker = SemanticChunker(target_tokens=50, overlap_tokens=10)
    chunks = chunker.chunk(doc, "test-file-id")

    assert len(chunks) >= 2
    assert all(c.chunk_id.startswith("test-file-id-chunk-") for c in chunks)
    assert any(c.heading_context for c in chunks)


def test_ner_and_keyword_extraction():
    sample_text = (
        "Satya Nadella announced Microsoft earnings on August 20, 2026 in Redmond. "
        "Total net revenue reached $56 billion dollars."
    )
    ner = EntityRecognizer()
    entities = ner.extract_entities(sample_text)
    assert len(entities) >= 1

    kw = KeywordExtractor()
    keywords = kw.extract_keywords(sample_text, top_n=5)
    assert len(keywords) >= 1
